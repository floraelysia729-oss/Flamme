"""
LLM Wiki 摄入管道 — PPT/PDF → Markdown

用法:
  python scripts/ingest.py <file_or_dir> [--level lite|pro] [--name 名称] [--ppt2pdf]

示例:
  python scripts/ingest.py "E:/课件/矩阵论.pptx" --level pro
  python scripts/ingest.py "D:/notebook/pro/数字系统设计" --level pro
  python scripts/ingest.py chap1.pdf --level lite --name "计算机系统"
"""

import argparse
import os
import re
import sys
import tempfile
import time
from pathlib import Path

# 确保 UTF-8 输出 + 禁用缓冲（管道/后台时也能看到进度）
sys.stdout.reconfigure(encoding="utf-8", line_buffering=True)

from flamme_paths import VAULT, converted_dir, ocr_dir


# ── PPT → PDF (Windows, via PowerPoint COM) ──────────────────────────

def pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    """用 PowerPoint COM 接口将 .pptx 转为 .pdf"""
    import comtypes.client

    pdf_path = out_dir / (pptx_path.stem + ".pdf")
    if pdf_path.exists():
        print(f"  [skip] PDF already exists: {pdf_path.name}")
        return pdf_path

    print(f"  [ppt→pdf] {pptx_path.name} ...")
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    abs_pptx = str(pptx_path.resolve())
    abs_pdf = str(pdf_path.resolve())
    deck = powerpoint.Presentations.Open(abs_pptx)
    deck.SaveAs(abs_pdf, 32)  # 32 = ppSaveAsPDF
    deck.Close()
    powerpoint.Quit()
    print(f"  [done] → {pdf_path.name}")
    return pdf_path


# ── PDF → Markdown (pymupdf4llm) ────────────────────────────────────

def pdf_to_markdown(pdf_path: Path) -> str:
    """用 pymupdf4llm 提取 PDF 文字为 markdown"""
    import pymupdf4llm

    print(f"  [extract] {pdf_path.name} ...")
    text = pymupdf4llm.to_markdown(str(pdf_path))
    return text


# ── PPT → Markdown (python-pptx, 直接提取) ──────────────────────────

def _table_to_md(table) -> str:
    """将 python-pptx Table 转为 Markdown 表格"""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if len(rows) >= 1:
        # 用第一行做表头
        header = rows[0]
        sep = "|" + "|".join("---" for _ in table.columns) + "|"
        return header + "\n" + sep + "\n" + "\n".join(rows[1:])
    return ""


def _format_text_frame(text_frame) -> str:
    """提取 text_frame 内容，保留粗体和列表层级"""
    lines = []
    for para in text_frame.paragraphs:
        if not para.text.strip():
            continue
        # 检测列表层级
        level = para.level
        prefix = "  " * level + "- " if level > 0 else "- "
        # 检测粗体 run
        runs_text = []
        for run in para.runs:
            if run.font.bold:
                runs_text.append(f"**{run.text}**")
            else:
                runs_text.append(run.text)
        line = "".join(runs_text).strip()
        if line:
            # 文本已有列表标记时只加缩进，避免 "- - text" 或 "- • text"
            if line.startswith(('- ', '• ', '* ')):
                indent = "  " * level
                lines.append(f"{indent}{line}")
            else:
                lines.append(f"{prefix}{line}")
    return "\n".join(lines)


def _extract_pptx_images(slide, tmp_dir: Path) -> list[Path]:
    """提取幻灯片中的嵌入图片，返回临时文件路径列表"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    images = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            ext = img.content_type.split("/")[-1]  # png, jpeg, etc.
            if ext == "jpeg":
                ext = "jpg"
            tmp_path = tmp_dir / f"slide_img_{len(images)}.{ext}"
            tmp_path.write_bytes(img.blob)
            images.append(tmp_path)
    return images


def pptx_to_markdown(pptx_path: Path, ocr_images: bool = False, ocr_interval: float = None) -> str:
    """用 python-pptx 直接提取 PPT 文字，支持表格和嵌入图片 OCR"""
    from pptx import Presentation

    print(f"  [extract] {pptx_path.name} ...")
    prs = Presentation(str(pptx_path))
    lines = [f"# {pptx_path.stem}\n"]

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)

        for i, slide in enumerate(prs.slides, 1):
            slide_title = ""
            body_parts = []
            notes_text = ""

            for shape in slide.shapes:
                # 文本框
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if not text:
                        continue
                    if not slide_title:
                        slide_title = text
                    else:
                        body_parts.append(_format_text_frame(shape.text_frame))

                # 表格
                elif shape.has_table:
                    md_table = _table_to_md(shape.table)
                    if md_table:
                        body_parts.append(md_table)

            # 提取嵌入图片 + OCR
            if ocr_images:
                imgs = _extract_pptx_images(slide, tmp_dir)
                if imgs:
                    print(f"    [ocr] slide {i}: {len(imgs)} image(s)")
                    ocr_results = ocr_images_files(imgs, min_interval=ocr_interval)
                    for img_path, ocr_text in zip(imgs, ocr_results):
                        if ocr_text:
                            body_parts.append(f"> [图片文字] {ocr_text}")

            # 备注
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_frame.text.strip()

            if slide_title or body_parts:
                lines.append(f"\n## Slide {i}: {slide_title}\n")
                for part in body_parts:
                    lines.append(part)
                    lines.append("")
                if notes_text:
                    lines.append(f"> **备注**: {notes_text}")

    return "\n".join(lines)


# ── 薄页检测 ────────────────────────────────────────────────────────

def detect_thin_pages(pdf_path: Path, threshold: int = 50) -> list[int]:
    """检测文字量低于阈值的页面（大概率是图表）"""
    import fitz

    doc = fitz.open(str(pdf_path))
    thin = []
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        # 去掉纯空白和标点后算字数
        clean = re.sub(r"\s+|[^\w]", "", text)
        if len(clean) < threshold:
            thin.append(i)
    doc.close()
    return thin


# ── 导出薄页图片 ───────────────────────────────────────────────────

def export_thin_pages(pdf_path: Path, pages: list[int], out_dir: Path) -> list[Path]:
    """将指定页面导出为 PNG"""
    import fitz

    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(pdf_path))
    exported = []
    for p in pages:
        if p - 1 < doc.page_count:
            page = doc[p - 1]
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            img_path = out_dir / f"{pdf_path.stem}_p{p}.png"
            pix.save(str(img_path))
            exported.append(img_path)
    doc.close()
    return exported


# ── OCR 识别 (RapidOCR HTTP API) ──────────────────────────────────

# 限速：两次 OCR 请求之间的最小间隔（秒），可通过 --rate-limit 覆盖
_OCR_MIN_INTERVAL = 1.0
_last_ocr_time = 0.0


def _ocr_single(img_path: Path, api_url: str = "http://localhost:1224/api/ocr",
                 max_retries: int = 3, min_interval: float = None) -> str | None:
    """OCR 单张图片，带限速和重试退避"""
    import base64
    import json
    import urllib.error
    import urllib.request

    global _last_ocr_time

    b64 = base64.b64encode(img_path.read_bytes()).decode()
    payload = json.dumps({"base64": b64}).encode()

    for attempt in range(max_retries):
        # 限速：确保请求间隔
        interval = min_interval if min_interval is not None else _OCR_MIN_INTERVAL
        elapsed = time.time() - _last_ocr_time
        if elapsed < interval:
            time.sleep(interval - elapsed)

        req = urllib.request.Request(api_url, data=payload, headers={"Content-Type": "application/json"})
        try:
            _last_ocr_time = time.time()
            with urllib.request.urlopen(req, timeout=30) as resp:
                data = json.loads(resp.read().decode())
                if data.get("code") in (100, 200) and data.get("data"):
                    texts = [item["text"] for item in data["data"] if isinstance(item, dict) and item.get("text")]
                    return "\n".join(texts).strip() or None
                return None
        except urllib.error.HTTPError as e:
            if e.code == 429:
                wait = 2 ** attempt * 2  # 2s, 4s, 8s
                print(f"    [ocr 429] {img_path.name}: rate limited, retry in {wait}s ({attempt+1}/{max_retries})")
                time.sleep(wait)
                continue
            print(f"    [ocr fail] {img_path.name}: HTTP {e.code}")
            return None
        except Exception as e:
            if attempt < max_retries - 1:
                wait = 2 ** attempt
                print(f"    [ocr retry] {img_path.name}: {e}, retry in {wait}s")
                time.sleep(wait)
                continue
            print(f"    [ocr fail] {img_path.name}: {e}")
            return None

    print(f"    [ocr fail] {img_path.name}: max retries exceeded")
    return None


def ocr_images_files(img_paths: list[Path], api_url: str = "http://localhost:1224/api/ocr",
                     min_interval: float = None) -> list[str]:
    """OCR 多张图片，返回对应文字列表（保持顺序，无文字为空字符串）"""
    results = []
    for i, img_path in enumerate(img_paths):
        text = _ocr_single(img_path, api_url, min_interval=min_interval)
        results.append(text or "")
        if (i + 1) % 10 == 0:
            print(f"    [ocr] {i+1}/{len(img_paths)} done")
    return results


def ocr_images(img_paths: list[Path], api_url: str = "http://localhost:1224/api/ocr",
               min_interval: float = None) -> dict[str, str]:
    """用 RapidOCR HTTP API 识别图片，返回 {文件名: 识别文字}"""
    results = {}
    for img_path in img_paths:
        text = _ocr_single(img_path, api_url, min_interval=min_interval)
        if text:
            results[img_path.name] = text
    return results


def ocr_thin_pages(pdf_path: Path, pages: list[int], source_dir: Path) -> str:
    """导出薄页图片 → umi-ocr 识别 → 返回合并的 OCR 文本"""
    imgs = export_thin_pages(pdf_path, pages, ocr_dir(source_dir))
    if not imgs:
        return ""
    print(f"  [ocr] recognizing {len(imgs)} pages via umi-ocr ...")
    ocr_results = ocr_images(imgs)
    if not ocr_results:
        print(f"  [ocr] no text recognized")
        return ""

    # 按页码排序拼成 markdown
    lines = []
    for img_path in sorted(imgs):
        text = ocr_results.get(img_path.name, "")
        if text:
            page_num = img_path.stem.split("_p")[-1]
            lines.append(f"\n### Page {page_num} (OCR)\n\n{text}\n")

    print(f"  [ocr] recognized {len(ocr_results)}/{len(imgs)} pages")
    return "\n".join(lines)


# ── Frontmatter ────────────────────────────────────────────────────

def make_frontmatter(title: str, source: str, level: str, tags: list[str] = None) -> str:
    from datetime import date

    today = date.today().isoformat()
    tag_str = str(tags or []).replace("'", '"')
    return f"""---
title: "{title}"
date: {today}
source: "{source}"
level: {level}
tags: {tag_str}
related: []
status: stable
---

"""


# ── 主流程 ─────────────────────────────────────────────────────────

def process_file(filepath: Path, level: str, name: str = None, ppt2pdf_flag: bool = False,
                  ocr_flag: bool = False, ocr_interval: float = None):
    """处理单个文件"""
    ext = filepath.suffix.lower()
    dir_name = name or filepath.stem
    source_dir = VAULT / level / dir_name
    source_dir.mkdir(parents=True, exist_ok=True)

    # 输出文件名用原始文件名，不用 --name
    file_stem = filepath.stem

    if ext == ".pdf":
        md_text = pdf_to_markdown(filepath)
        source = str(filepath.relative_to(VAULT)) if filepath.is_relative_to(VAULT) else filepath.name

        # 薄页 OCR
        thin = detect_thin_pages(filepath)
        ocr_text = ""
        if thin:
            print(f"  [thin pages] {len(thin)} pages: {thin}")
            if ocr_flag:
                ocr_text = ocr_thin_pages(filepath, thin, source_dir)
            else:
                imgs = export_thin_pages(filepath, thin, ocr_dir(source_dir))
                print(f"  → exported {len(imgs)} images (use --ocr to auto-recognize)")
        else:
            print("  [thin pages] none")

        # 合并：提取文字 + OCR 补充
        if ocr_text:
            md_text = f"{md_text}\n\n---\n\n# OCR 补充内容（薄页）\n{ocr_text}"

        md_text = make_frontmatter(file_stem, source, level) + md_text
        out_md = converted_dir(source_dir) / f"{file_stem}.md"
        out_md.write_text(md_text, encoding="utf-8")
        print(f"  → {out_md}")

        return out_md

    elif ext == ".pptx":
        # 方式1: 直接用 python-pptx 提取文字 + 可选图片 OCR
        md_text = pptx_to_markdown(filepath, ocr_images=ocr_flag, ocr_interval=ocr_interval)
        source = filepath.name

        # 方式2: 如果指定 --ppt2pdf，额外生成 PDF 并提取
        if ppt2pdf_flag:
            try:
                pdf_path = pptx_to_pdf(filepath, converted_dir(source_dir))
                pdf_text = pdf_to_markdown(pdf_path)
                # 合并：PPTX 提取的备注 + PDF 提取的视觉文字
                md_text = f"{md_text}\n\n---\n\n# PDF 提取补充内容\n\n{pdf_text}"
                source = f"{filepath.name} (pptx+pdf)"

                # 薄页检测
                thin = detect_thin_pages(pdf_path)
                if thin:
                    print(f"  [thin pages] {len(thin)} pages: {thin}")
                    imgs = export_thin_pages(pdf_path, thin, ocr_dir(source_dir))
                    print(f"  → exported {len(imgs)} images")
            except Exception as e:
                print(f"  [warn] ppt2pdf failed: {e}")

        md_text = make_frontmatter(file_stem, source, level) + md_text
        out_md = converted_dir(source_dir) / f"{file_stem}.md"
        out_md.write_text(md_text, encoding="utf-8")
        print(f"  → {out_md}")
        return out_md

    elif ext == ".md":
        # 已有 markdown，只复制
        import shutil
        out_md = converted_dir(source_dir) / filepath.name
        shutil.copy2(filepath, out_md)
        print(f"  → copied to {out_md}")
        return out_md

    else:
        print(f"  [skip] unsupported format: {ext}")
        return None


def main():
    parser = argparse.ArgumentParser(description="LLM Wiki 摄入管道")
    parser.add_argument("path", help="文件或目录路径")
    parser.add_argument("--level", choices=["raw", "lite", "pro"], default="lite", help="处理级别")
    parser.add_argument("--name", help="自定义目录名（默认用文件名）")
    parser.add_argument("--ppt2pdf", action="store_true", help="PPT 额外转 PDF 提取（需安装 PowerPoint）")
    parser.add_argument("--ocr", action="store_true", help="用 umi-ocr 自动识别薄页图片（需 umi-ocr 运行中）")
    parser.add_argument("--rate-limit", type=float, default=1.0, help="OCR 请求间隔秒数（默认1.0，防限流）")
    args = parser.parse_args()

    path = Path(args.path)
    if not path.exists():
        print(f"Error: {path} not found")
        sys.exit(1)

    # 收集要处理的文件
    supported = {".pdf", ".pptx"}
    if path.is_dir():
        files = sorted(p for p in path.rglob("*")
                       if p.suffix.lower() in supported and not p.name.startswith("~$"))
    else:
        files = [path]

    if not files:
        print("No supported files found")
        sys.exit(1)

    print(f"Found {len(files)} file(s) to process (level={args.level})\n")

    results = []
    for f in files:
        print(f"[*] {f.name}")
        result = process_file(f, args.level, args.name, args.ppt2pdf, args.ocr, args.rate_limit)
        if result:
            results.append(result)
        print()

    print(f"Done. {len(results)} file(s) processed.")
    if results:
        print("Output:")
        for r in results:
            print(f"  {r}")


if __name__ == "__main__":
    main()
